"""
Builds the slide deck `outputs/Company_A_Churn_Deck.pptx` from results.json and
the figures in outputs/figures/. <= 15 content slides.

Run:  python build_deck.py
"""
from __future__ import annotations

import json
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Emu, Inches, Pt

ROOT = Path(__file__).resolve().parent
FIG = ROOT / "outputs" / "figures"
RES = json.loads((ROOT / "outputs" / "results" / "results.json").read_text())

# ---- palette (3 brand + neutrals) -----------------------------------------
INK = RGBColor(0x14, 0x1B, 0x2E)      # near-black navy (text)
NAVY = RGBColor(0x1F, 0x3A, 0x5F)     # primary brand
TEAL = RGBColor(0x12, 0x8A, 0x86)     # accent / positive
AMBER = RGBColor(0xC9, 0x7A, 0x1E)    # accent / caution
SLATE = RGBColor(0x5B, 0x66, 0x77)    # muted text
CLOUD = RGBColor(0xF2, 0xF4, 0xF7)    # light panel
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
RED = RGBColor(0xB1, 0x3A, 0x3A)      # negative

W, H = Inches(13.333), Inches(7.5)    # 16:9

prs = Presentation()
prs.slide_width = W
prs.slide_height = H
BLANK = prs.slide_layouts[6]


def slide():
    return prs.slides.add_slide(BLANK)


def rect(s, x, y, w, h, color):
    from pptx.enum.shapes import MSO_SHAPE
    sp = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    sp.fill.solid()
    sp.fill.fore_color.rgb = color
    sp.line.fill.background()
    sp.shadow.inherit = False
    return sp


def text(s, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
         space_after=6, line_spacing=1.12):
    """runs: list of (string, size, color, bold) OR list of such lists (paragraphs)."""
    tb = s.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    if runs and not isinstance(runs[0], list):
        runs = [runs]
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space_after)
        p.line_spacing = line_spacing
        for (txt, size, color, bold) in para:
            r = p.add_run()
            r.text = txt
            r.font.size = Pt(size)
            r.font.color.rgb = color
            r.font.bold = bold
            r.font.name = "Calibri"
    return tb


def header(s, kicker, title):
    rect(s, 0, 0, W, Inches(0.18), TEAL)
    text(s, Inches(0.6), Inches(0.42), Inches(12), Inches(0.4),
         [(kicker.upper(), 12, TEAL, True)])
    text(s, Inches(0.6), Inches(0.72), Inches(12.1), Inches(0.9),
         [(title, 28, INK, True)])
    rect(s, Inches(0.6), Inches(1.62), Inches(1.1), Inches(0.045), AMBER)


def img(s, name, x, y, w, h=None):
    p = FIG / name
    kw = {"width": w} if h is None else {"width": w, "height": h}
    s.shapes.add_picture(str(p), x, y, **kw)


def fnum(x):
    return f"${x:,.0f}"


# ============================================================ 1. TITLE
s = slide()
rect(s, 0, 0, W, H, NAVY)
rect(s, 0, Inches(5.0), W, Inches(0.07), TEAL)
text(s, Inches(0.9), Inches(1.7), Inches(11.5), Inches(2),
     [[("Reducing Customer Churn at Company A", 40, WHITE, True)]])
text(s, Inches(0.9), Inches(3.1), Inches(11.5), Inches(1.4),
     [[("A profit-driven, calibrated, and fair churn model — ", 20, CLOUD, False),
       ("predict, prioritize, and act.", 20, RGBColor(0x9F, 0xD8, 0xD6), True)]])
text(s, Inches(0.9), Inches(5.3), Inches(11.5), Inches(1.4),
     [[("Telecom Retention Analytics", 14, RGBColor(0xB9, 0xC4, 0xD4), True)],
      [("100,000 customers  |  123 engineered behavioral features  |  "
        "demographics excluded by design", 13, RGBColor(0x9AA, 0x0, 0x0) if False else RGBColor(0xB9, 0xC4, 0xD4), False)]])

# ============================================================ 2. AGENDA / TL;DR
s = slide()
header(s, "Executive summary", "The bottom line")
d = RES
base = next(x for x in d["scenarios"] if x["name"] == "Base")
panels = [
    ("Realistic model", f"AUC {d['calibration']['auc_calibrated']:.2f}",
     "Tuned + calibrated XGBoost. No leakage, no 0.99 fantasy.", TEAL),
    ("Targeting lift", f"{d['capture_at_20pct']['lift_evar_vs_random']:.2f}x",
     "vs. random outreach, ranking by expected value at risk.", NAVY),
    ("Base-case ROI", f"{base['roi']:.1f}x",
     f"Net {fnum(base['net_profit_population'])} at the profit-optimal "
     f"{int(d['profit_curve_base']['optimal_frac']*100)}% contact depth.", AMBER),
]
x = Inches(0.6)
for (k, big, sub, col) in panels:
    rect(s, x, Inches(2.1), Inches(3.9), Inches(2.6), CLOUD)
    rect(s, x, Inches(2.1), Inches(3.9), Inches(0.12), col)
    text(s, x + Inches(0.3), Inches(2.35), Inches(3.3), Inches(0.4),
         [(k.upper(), 12, col, True)])
    text(s, x + Inches(0.3), Inches(2.75), Inches(3.3), Inches(0.9),
         [(big, 40, INK, True)])
    text(s, x + Inches(0.3), Inches(3.75), Inches(3.4), Inches(0.9),
         [(sub, 13, SLATE, False)], line_spacing=1.15)
    x += Inches(4.25)
text(s, Inches(0.6), Inches(5.1), Inches(12.1), Inches(1.8),
     [[("What changed from a naive first pass:  ", 14, INK, True),
       ("protected attributes removed, probabilities calibrated, scores "
        "re-weighted to the true ~2% churn rate, the decision threshold set on a "
        "profit curve, drivers explained with SHAP, and the business case shown "
        "as bear / base / bull with cited assumptions.", 14, SLATE, False)]],
     line_spacing=1.25)

# ============================================================ 3. PROBLEM & DATA
s = slide()
header(s, "Problem framing", "Churn is expensive — but not every customer is worth the same call")
text(s, Inches(0.6), Inches(2.0), Inches(6.0), Inches(4.5),
     [[("The task.", 16, NAVY, True)],
      [("Predict which subscribers will churn, explain why, and turn that into a "
        "retention plan that makes money.", 14, SLATE, False)],
      [("", 6, SLATE, False)],
      [("The data.", 16, NAVY, True)],
      [(f"{d['data']['n_rows']:,} customers, {d['data']['n_raw_cols']} raw columns "
        "across account, usage, and demographic tables.", 14, SLATE, False)],
      [("", 6, SLATE, False)],
      [("The honesty problem.", 16, AMBER, True)],
      [(f"The sample is balanced to ~50% churn for training, but real monthly "
        f"churn is ~2%. We train on the sample, then re-weight predictions back "
        f"to reality so the dollars are credible.", 14, SLATE, False)]],
     line_spacing=1.2)
rect(s, Inches(7.0), Inches(2.0), Inches(5.7), Inches(4.3), CLOUD)
text(s, Inches(7.3), Inches(2.25), Inches(5.1), Inches(0.5),
     [("KEY NUMBERS", 12, TEAL, True)])
kv = [
    ("ARPU (avg monthly revenue)", f"${d['descriptive']['arpu']:.0f}"),
    ("Median tenure", f"{d['descriptive']['median_tenure_months']:.0f} months"),
    ("Model features (engineered)", f"{d['data']['n_model_features']}"),
    ("Protected attributes dropped", f"{len(d['data']['dropped_protected'])}"),
    ("Assumed true monthly churn", f"{d['base_rate']['true_rate']*100:.0f}%"),
]
yy = Inches(2.8)
for (k, v) in kv:
    text(s, Inches(7.3), yy, Inches(3.6), Inches(0.5), [(k, 13, INK, False)])
    text(s, Inches(10.9), yy, Inches(1.5), Inches(0.5), [(v, 15, NAVY, True)],
         align=PP_ALIGN.RIGHT)
    yy += Inches(0.66)

# ============================================================ 4. FAIRNESS
s = slide()
header(s, "Responsible modeling", "We deliberately excluded demographics")
text(s, Inches(0.6), Inches(2.0), Inches(12.1), Inches(1.2),
     [[("Targeting retention by ethnicity, income, marital status, or dwelling is "
        "a legal and ethical hazard — and it isn't actionable. ", 15, SLATE, False),
       ("All 20 such attributes are dropped before the model sees the data, "
        "with an automated guard that fails the build if one leaks through "
        "encoding.", 15, INK, True)]], line_spacing=1.25)
cols = [d["data"]["dropped_protected"][i::4] for i in range(4)]
x = Inches(0.6)
for c in cols:
    rect(s, x, Inches(3.5), Inches(2.95), Inches(3.0), CLOUD)
    body = [[("EXCLUDED", 11, RED, True)]]
    body += [[("• " + a, 13, SLATE, False)] for a in c]
    text(s, x + Inches(0.25), Inches(3.7), Inches(2.6), Inches(2.7), body,
         line_spacing=1.2)
    x += Inches(3.05)

# ============================================================ 5. METHOD PIPELINE
s = slide()
header(s, "Approach", "One reproducible pipeline, end to end")
steps = [
    ("1  Engineer", "Merge usage + account tables; build behavioral features "
     "(equipment age, usage trend, care intensity, overage share)."),
    ("2  Model", "Compare Logistic Regression, Random Forest, XGBoost; tune XGBoost "
     "with randomized search + 5-fold CV."),
    ("3  Calibrate", "Isotonic calibration so scores are real probabilities; "
     "re-weight to the true 2% base rate."),
    ("4  Decide", "Rank by expected value at risk; choose contact depth on the "
     "profit curve, not an arbitrary 0.5 cutoff."),
    ("5  Explain & act", "SHAP drivers, K-means personas, and an uplift/persuadables "
     "view feed a bear/base/bull business case."),
]
yy = Inches(2.0)
for (t, b) in steps:
    rect(s, Inches(0.6), yy, Inches(2.5), Inches(0.82), NAVY)
    text(s, Inches(0.7), yy + Inches(0.12), Inches(2.3), Inches(0.6),
         [(t, 16, WHITE, True)], anchor=MSO_ANCHOR.MIDDLE)
    text(s, Inches(3.35), yy + Inches(0.05), Inches(9.3), Inches(0.78),
         [(b, 14, SLATE, False)], anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.1)
    yy += Inches(0.95)

# ============================================================ 6. MODEL COMPARISON
s = slide()
header(s, "Model selection", "XGBoost wins — at a believable accuracy")
img(s, "model_comparison.png", Inches(0.6), Inches(1.95), Inches(7.4))
xg = d["model_comparison"]["XGBoost"]
text(s, Inches(8.3), Inches(2.1), Inches(4.4), Inches(4.5),
     [[("Why this is credible", 16, NAVY, True)],
      [(f"AUC {xg['auc']:.2f}, PR-AUC {xg['pr_auc']:.2f}. ", 14, INK, True),
       ("An AUC near 0.70 is normal for telecom churn on behavioral data — "
        "a 0.95+ score would signal target leakage.", 14, SLATE, False)],
      [("", 6, SLATE, False)],
      [(f"5-fold CV: {d['cv']['mean_auc']:.3f} ± {d['cv']['std_auc']:.3f}. ",
        14, INK, True),
       ("Tiny variance across folds means the model is stable, not lucky.",
        14, SLATE, False)],
      [("", 6, SLATE, False)],
      [("Brier score improves after calibration — the probabilities can be "
        "trusted for dollar math.", 14, SLATE, False)]],
     line_spacing=1.2)

# ============================================================ 7. CALIBRATION
s = slide()
header(s, "Trustworthy probabilities", "Calibration + the base-rate correction")
img(s, "calibration.png", Inches(0.6), Inches(1.95), Inches(7.0))
br = d["base_rate"]
text(s, Inches(7.9), Inches(2.1), Inches(4.8), Inches(4.5),
     [[("A score is not a probability", 16, NAVY, True)],
      [("Raw model scores cluster near 0.5 because the training sample is "
        "balanced. Isotonic calibration aligns predicted probability with "
        "observed frequency (the diagonal).", 14, SLATE, False)],
      [("", 6, SLATE, False)],
      [("Then we correct the prior:", 15, AMBER, True)],
      [(f"mean P(churn) {br['mean_p_sample']:.2f} in the sample  →  "
        f"{br['mean_p_true']:.3f} re-weighted to the real ~2% rate.",
        14, INK, True)],
      [("", 6, SLATE, False)],
      [("Without this step, every revenue-at-risk and ROI figure would be "
        "inflated by ~20x.", 14, SLATE, False)]],
     line_spacing=1.2)

# ============================================================ 8. DRIVERS (SHAP)
s = slide()
header(s, "Why customers leave", "SHAP — the model's honest drivers")
img(s, "shap_importance.png", Inches(0.6), Inches(1.95), Inches(7.4))
text(s, Inches(8.3), Inches(2.1), Inches(4.4), Inches(4.5),
     [[("All actionable, none demographic", 16, NAVY, True)],
      [("Equipment age (eqpdays)", 14, TEAL, True),
       (" — the top driver. Old handsets churn; proactive upgrade offers are the "
        "lever.", 14, SLATE, False)],
      [("Tenure & usage trend", 14, TEAL, True),
       (" — declining minutes-of-use flags disengagement early.", 14, SLATE, False)],
      [("Care intensity & dropped calls", 14, TEAL, True),
       (" — service quality and support friction. Fixable operationally.",
        14, SLATE, False)],
      [("", 6, SLATE, False)],
      [("Because demographics were excluded, every top driver maps to an action "
        "the business can actually take.", 14, INK, True)]],
     line_spacing=1.2)

# ============================================================ 9. HANDSET HONESTY
s = slide()
header(s, "Honest segmentation", "Churn rises with equipment age")
img(s, "handset_churn.png", Inches(0.6), Inches(1.95), Inches(7.4))
text(s, Inches(8.3), Inches(2.1), Inches(4.4), Inches(4.5),
     [[("Binned, not over-precise", 16, NAVY, True)],
      [("Equipment age is grouped into interpretable buckets rather than implying "
        "false day-level precision.", 14, SLATE, False)],
      [("", 6, SLATE, False)],
      [("The pattern is monotonic and clean: the longer a customer keeps an aging "
        "device, the more likely they leave.", 14, INK, True)],
      [("", 6, SLATE, False)],
      [("Action: ", 14, AMBER, True),
       ("trigger an upgrade offer as devices cross the high-risk age band.",
        14, SLATE, False)]],
     line_spacing=1.2)

# ============================================================ 10. WHO TO CALL (capture)
s = slide()
header(s, "Prioritize", "Call the highest-value risk first")
img(s, "capture_curve.png", Inches(0.6), Inches(1.95), Inches(7.4))
cap = d["capture_at_20pct"]
text(s, Inches(8.3), Inches(2.1), Inches(4.4), Inches(4.5),
     [[("Expected value at risk", 16, NAVY, True)],
      [("We rank by ", 14, SLATE, False),
       ("P(churn) × margin × lifetime value", 14, INK, True),
       (", not probability alone — a $100 customer at 30% risk beats a $20 "
        "customer at 60%.", 14, SLATE, False)],
      [("", 6, SLATE, False)],
      [(f"Contacting the top 20% by value captures "
        f"{cap['by_evar']*100:.0f}% of churn value — "
        f"{cap['lift_evar_vs_random']:.2f}x better than random.", 14, TEAL, True)]],
     line_spacing=1.2)

# ============================================================ 11. HOW MANY (profit)
s = slide()
header(s, "Optimize", "How deep to go — the profit curve")
img(s, "profit_curve.png", Inches(0.6), Inches(1.95), Inches(7.4))
pc = d["profit_curve_base"]
text(s, Inches(8.3), Inches(2.1), Inches(4.4), Inches(4.5),
     [[("Maximize profit, not recall", 16, NAVY, True)],
      [("Calling everyone wastes incentives; calling too few leaves value on the "
        "table. We sweep contact depth and pick the peak.", 14, SLATE, False)],
      [("", 6, SLATE, False)],
      [(f"Optimal depth: {int(pc['optimal_frac']*100)}%", 16, AMBER, True)],
      [(f"Net profit {fnum(pc['optimal_net_profit_population'])}  •  "
        f"ROI {pc['optimal_roi']:.1f}x  •  "
        f"~{pc['customers_saved_population']:.0f} customers saved (base case).",
        14, INK, True)]],
     line_spacing=1.2)

# ============================================================ 12. PERSONAS
s = slide()
header(s, "Who they are", "Four behavioral personas (K-means)")
img(s, "personas.png", Inches(0.6), Inches(1.95), Inches(7.2))
# pick the standout persona
cl = sorted(d["personas"]["clusters"], key=lambda c: -c["churn_rate"])[0]
text(s, Inches(8.1), Inches(2.1), Inches(4.6), Inches(4.5),
     [[("Behavioral, not demographic", 16, NAVY, True)],
      [("Clusters are built only from behavior — usage trend, equipment age, care "
        "intensity, overage, tenure.", 14, SLATE, False)],
      [("", 6, SLATE, False)],
      [("High-risk segment", 14, RED, True),
       (f": churn rate {cl['churn_rate']*100:.0f}%, high care-call intensity and "
        f"aging equipment — the clearest save target.", 14, SLATE, False)],
      [("", 6, SLATE, False)],
      [("A separate high-ARPU segment justifies premium, concierge-style retention "
        "offers.", 14, INK, True)]],
     line_spacing=1.2)

# ============================================================ 13. UPLIFT / PERSUADABLES
s = slide()
header(s, "Persuadables", "Treat customers a campaign can actually move")
img(s, "uplift.png", Inches(0.6), Inches(1.95), Inches(7.2))
up = d["uplift"]
text(s, Inches(8.1), Inches(2.05), Inches(4.6), Inches(4.7),
     [[("Beyond risk: responsiveness", 16, NAVY, True)],
      [("Some high-risk customers leave regardless; the money is in ", 14, SLATE, False),
       ("persuadables", 14, TEAL, True),
       (" — those whose retention an offer actually changes (T-learner uplift).",
        14, SLATE, False)],
      [("", 5, SLATE, False)],
      [("Honesty note", 13, AMBER, True)],
      [("This data is observational (no test/control), so the response here is a "
        "clearly-labeled simulation to demonstrate the method.", 13, SLATE, False)],
      [("", 5, SLATE, False)],
      [("Next step: a live A/B holdout to measure real uplift before scaling.",
        13, INK, True)]],
     line_spacing=1.12)

# ============================================================ 14. BUSINESS CASE
s = slide()
header(s, "The money", "Bear / base / bull — profit-positive in the base & bull")
img(s, "scenarios.png", Inches(0.6), Inches(1.95), Inches(6.8))
text(s, Inches(7.7), Inches(2.0), Inches(5.0), Inches(5.0),
     [[("Assumptions are cited, not invented", 15, NAVY, True)],
      [("Each scenario varies save rate, offer cost, margin, and CLV horizon — "
        "all sourced in REFERENCES.md.", 13, SLATE, False)]],
     line_spacing=1.12)
yy = Inches(3.1)
for sc in d["scenarios"]:
    col = {"Bear": RED, "Base": NAVY, "Bull": TEAL}[sc["name"]]
    np_pop = sc["net_profit_population"]
    rect(s, Inches(7.7), yy, Inches(5.0), Inches(1.05), CLOUD)
    rect(s, Inches(7.7), yy, Inches(0.12), Inches(1.05), col)
    text(s, Inches(7.95), yy + Inches(0.1), Inches(2.2), Inches(0.85),
         [[(sc["name"], 16, col, True)],
          [(f"save {int(sc['save_rate']*100)}% · ${sc['offer_cost']:.0f} offer",
            11, SLATE, False)]], line_spacing=1.0)
    text(s, Inches(10.1), yy + Inches(0.1), Inches(2.5), Inches(0.85),
         [[(fnum(np_pop), 19, (RED if np_pop < 0 else INK), True)],
          [(f"ROI {sc['roi']:.1f}x · ~{sc['customers_saved_population']:.0f} saved",
            11, SLATE, False)]], align=PP_ALIGN.RIGHT, line_spacing=1.0)
    yy += Inches(1.15)

# ============================================================ 15. RECOMMENDATIONS
s = slide()
rect(s, 0, 0, W, H, NAVY)
rect(s, 0, 0, W, Inches(0.18), TEAL)
text(s, Inches(0.7), Inches(0.5), Inches(12), Inches(0.9),
     [("RECOMMENDATIONS & NEXT STEPS", 13, RGBColor(0x9F, 0xD8, 0xD6), True)])
text(s, Inches(0.7), Inches(0.85), Inches(12), Inches(0.9),
     [("Operationalize the model — and prove it with a test", 26, WHITE, True)])
recs = [
    ("Deploy value-ranked scoring", "Score the base monthly; route the profit-optimal "
     f"top {int(d['profit_curve_base']['optimal_frac']*100)}% by value-at-risk to retention."),
    ("Trigger device-age offers", "Automate upgrade offers as handsets cross the "
     "high-risk age band — the #1 SHAP driver."),
    ("Fix service friction", "Feed dropped-call and care-intensity signals to network "
     "and support ops as leading indicators."),
    ("Run an A/B retention test", "Randomized treatment/control to measure REAL uplift "
     "and replace the simulated persuadables figures."),
    ("Re-validate assumptions", "Confirm save rate, offer cost, margin, and CLV with "
     "finance; refresh the bear/base/bull case quarterly."),
]
yy = Inches(2.05)
for i, (t, b) in enumerate(recs, 1):
    rect(s, Inches(0.7), yy, Inches(0.5), Inches(0.5), TEAL)
    text(s, Inches(0.7), yy, Inches(0.5), Inches(0.5), [(str(i), 18, WHITE, True)],
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    text(s, Inches(1.45), yy - Inches(0.02), Inches(11), Inches(0.55),
         [[(t + ".  ", 16, WHITE, True), (b, 14, RGBColor(0xC8, 0xD2, 0xE0), False)]],
         anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.05)
    yy += Inches(0.84)
text(s, Inches(0.7), Inches(6.6), Inches(12), Inches(0.6),
     [[("Reproducible: ", 12, RGBColor(0x9F, 0xD8, 0xD6), True),
       ("churn_analysis.ipynb regenerates every figure and number from the raw "
        "CSVs; assumptions cited in REFERENCES.md.", 12, RGBColor(0xB9, 0xC4, 0xD4), False)]])

out = ROOT / "outputs" / "Company_A_Churn_Deck.pptx"
prs.save(str(out))
print(f"Saved {out} with {len(prs.slides._sldIdLst)} slides")
